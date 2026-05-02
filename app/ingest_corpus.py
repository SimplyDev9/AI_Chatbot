#app.ingest_corpus.py
import os
import time
import shutil
from pathlib import Path
from datetime import datetime
from dotenv import load_dotenv
from langchain_aws import BedrockEmbeddings
from langchain_chroma import Chroma
from langchain_text_splitters import RecursiveCharacterTextSplitter
import docx2txt
from PyPDF2 import PdfReader
from app.config import CHROMA_DIR
import hashlib
from app.document_loader import load_document
from app.logger import log_ingest
from pptx import Presentation
import pandas as pd
import pdfplumber
import re
from app.logger import logger
import concurrent.futures
import math
import time
from functools import lru_cache

load_dotenv()


@lru_cache(maxsize=2)
def get_vectordb_instance():
    embed_model = BedrockEmbeddings(
        model_id="cohere.embed-english-v3",
        region_name=os.getenv("AWS_DEFAULT_REGION"),
        aws_access_key_id=os.getenv("AWS_ACCESS_KEY_ID"),
        aws_secret_access_key=os.getenv("AWS_SECRET_ACCESS_KEY"),
        # model_kwargs={
        #     "dimensions": 512,
        #     "normalize": True
        # }
    )

    return Chroma(
        persist_directory=str(CHROMA_DIR),
        embedding_function=embed_model
    )


def safe_clear_chroma_db():
    """Safely clears the Chroma DB directory (handles Windows file locks)."""
    if not os.path.exists(CHROMA_DIR):
        log_ingest(message="no_chroma_db", action="safe_clear", status="skipped", dir=str(CHROMA_DIR))
        return

    log_ingest(message="attempt_clear_old_chroma", action="safe_clear", status="started", dir=str(CHROMA_DIR))
    max_retries = 5
    for attempt in range(max_retries):
        try:
            shutil.rmtree(CHROMA_DIR)
            log_ingest(message="cleared_old_chroma", action="safe_clear", status="success", dir=str(CHROMA_DIR))
            return
        except PermissionError as e:
            log_ingest(message="chroma_locked", action="safe_clear", status="locked", attempt=attempt + 1, error=str(e))
            time.sleep(1)
        except Exception as e:
            log_ingest(message="clear_chroma_failed", action="safe_clear", status="failed", error=str(e))
            break
    log_ingest(message="clear_chroma_retry_failed", action="safe_clear", status="failed", attempts=max_retries)


def extract_text(fp: Path) -> str:
    ext = fp.suffix.lower()
    text = ""

    try:
        if ext == ".txt":
            text = fp.read_text(encoding="utf-8", errors="ignore")

        elif ext == ".docx":
            text = docx2txt.process(str(fp))

        elif ext == ".pdf":
            pages_text = []
            with pdfplumber.open(str(fp)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text() or ""
                    tables = page.extract_tables()
                    for table in tables:
                        for row in table:
                            cleaned_row = [str(cell).strip() if cell is not None else "" for cell in row]
                            page_text += "\n" + " | ".join(cleaned_row)
                    pages_text.append(page_text)
            text = ("\n".join(pages_text))

        elif ext == ".csv":
            df = pd.read_csv(fp, encoding="utf-8", encoding_errors="ignore")
            text = df.to_string()

        elif ext in [".xls", ".xlsx"]:
            df = pd.read_excel(fp)
            text = df.to_string()

        elif ext == ".pptx":
            presentation = Presentation(str(fp))
            slides_text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text.append(shape.text)
            text = "\n".join(slides_text)

        else:
            log_ingest(message="unsupported_file_type", action="extract_text", filename=fp.name, error=f"No extractor for {ext}")

    except Exception as e:
        log_ingest(message="failed_read", action="extract_text", filename=fp.name, error=str(e))

    return text.strip()

def ingest(corpus_dir=None, clear=False):
    if corpus_dir is None:
        corpus_dir = Path("Python_Project/corpus").resolve()
    else:
        corpus_dir = Path(corpus_dir).resolve()

    if not corpus_dir.exists():
        log_ingest(message="corpus_not_found", action="ingest", dir=str(corpus_dir))
        return

    # ✅ Use the safe clear function
    if clear:
        safe_clear_chroma_db()

    # Initialize embeddings and Chroma DB
    embed_model = BedrockEmbeddings(
        model_id="cohere.embed-english-v3",
        region_name=os.getenv("AWS_DEFAULT_REGION"),
        aws_access_key_id=os.getenv("AWS_ACCESS_KEY_ID"),
        aws_secret_access_key=os.getenv("AWS_SECRET_ACCESS_KEY"),
        # model_kwargs={
        #     "dimensions": 512,
        #     "normalize": True
        # }
    )

    vectordb = Chroma(
        persist_directory=str(CHROMA_DIR),
        embedding_function=embed_model
    )

    vectordb = Chroma(
        persist_directory=str(CHROMA_DIR),
        embedding_function=embed_model
    )

    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=150)
    file_paths = list(corpus_dir.glob("*"))

    for fpath in file_paths:

        file_hash = calculate_file_hash(fpath)

        # Check if file already exists
        existing = vectordb.get(where={"filename": fpath.name})

        if existing and existing.get("metadatas"):
            existing_hash = existing["metadatas"][0].get("file_hash")

            if existing_hash == file_hash:
                log_ingest(message="skip_unchanged", action="ingest", filename=fpath.name, status="skipped")
                continue

            else:
                log_ingest(message="file_changed", action="ingest", filename=fpath.name, status="updating")
                vectordb.delete(where={"filename": fpath.name})

        text = load_document(str(fpath))

        if not text:
            continue

        chunks = splitter.split_text(text)

        for i, chunk in enumerate(chunks, start=1):
            vectordb.add_texts(
                texts=[chunk],
                metadatas=[{
                    "filename": fpath.name,
                    "file_hash": file_hash,
                    "chunk_index": i,
                    "ingestion_date": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    "source": "local",
                    "file_type": fpath.suffix.lower(),
                    "file_size": os.path.getsize(fpath),
                }]
            )

        log_ingest(message="file_ingested", action="ingest", filename=fpath.name, chunks=len(chunks))


def calculate_file_hash(file_path: Path):
    """Generate hash for file to detect changes"""
    hasher = hashlib.md5()

    with open(file_path, "rb") as f:
        buf = f.read()
        hasher.update(buf)

    return hasher.hexdigest()



def ingest_single_file(file_path: Path, metadata=None):
    try:
        start_ts = time.time()
        logger.info("[INGEST] Starting ingestion for file: %s", file_path.name)
        logger.debug("[INGEST] File path: %s | Exists: %s | Size: %s",
                     str(file_path), file_path.exists(), file_path.stat().st_size if file_path.exists() else "n/a")
        if metadata:
            logger.debug("[INGEST] Incoming metadata: %s", metadata)

        vectordb = get_vectordb_instance()
        logger.debug("[INGEST] Obtained vectordb instance: %s", type(vectordb))

        file_hash = calculate_file_hash(file_path)
        logger.info("[INGEST] Calculated file hash: %s", file_hash)

        existing = vectordb.get(where={"filename": file_path.name})
        logger.debug("[INGEST] Existing lookup result: %s", existing)

        if existing and existing.get("metadatas"):
            existing_hash = existing["metadatas"][0].get("file_hash")
            logger.info("[INGEST] Found existing file entry: existing_hash=%s", existing_hash)
            if existing_hash == file_hash:
                logger.info("[INGEST] File unchanged, skipping ingestion: %s", file_path.name)
                logger.debug("[INGEST] Total time for skip: %.3fs", time.time() - start_ts)
                return
            logger.info("[INGEST] File changed, deleting old entries for: %s", file_path.name)
            try:
                vectordb.delete(where={"filename": file_path.name})
                logger.debug("[INGEST] Deleted old entries for: %s", file_path.name)
            except Exception as e:
                logger.warning("[INGEST] Failed to delete old entries for %s: %s", file_path.name, str(e))

        logger.info("[INGEST] Extracting text from file: %s", file_path.name)
        text = extract_text(file_path)
        if not text:
            logger.warning("[INGEST] No text extracted from file, aborting: %s", file_path.name)
            logger.debug("[INGEST] Total time when extraction empty: %.3fs", time.time() - start_ts)
            return
        logger.debug("[INGEST] Extracted text length: %d", len(text))

        # adaptive chunking
        if len(text) > 200_000:
            chunk_size = 800
            chunk_overlap = 100
            batch_size = 16
            logger.info("[INGEST] Using large-file chunking: chunk_size=%d, chunk_overlap=%d, batch_size=%d",
                        chunk_size, chunk_overlap, batch_size)
        else:
            chunk_size = int(os.getenv("CHUNK_SIZE", "300"))
            chunk_overlap = int(os.getenv("CHUNK_OVERLAP", "100"))
            batch_size = 10
            logger.info("[INGEST] Using default chunking: chunk_size=%d, chunk_overlap=%d, batch_size=%d",
                        chunk_size, chunk_overlap, batch_size)

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", "Q:", "A:", ".", " "]
        )

        chunks = splitter.split_text(text)
        logger.info("[INGEST] Split into %d chunks (approx).", len(chunks))

        metadatas = [
            {
                "filename": file_path.name,
                "file_hash": file_hash,
                "chunk_index": i,
                "file_type": file_path.suffix.lower(),
                "chunk_length": len(chunks[i]),
                "ingested_at": datetime.now().isoformat()
            }
            for i in range(len(chunks))
        ]

        logger.debug("[INGEST] Prepared metadatas for %d chunks.", len(metadatas))

        batches = [
            (chunks[i:i + batch_size], metadatas[i:i + batch_size])
            for i in range(0, len(chunks), batch_size)
        ]
        logger.info("[INGEST] Created %d batches (batch_size=%d).", len(batches), batch_size)

        for i, (text_batch, meta_batch) in enumerate(batches):
            batch_start = time.time()
            logger.info("[INGEST] Processing batch %d/%d with %d chunks", i + 1, len(batches), len(text_batch))
            success = False
            for attempt in range(5):
                try:
                    logger.debug("[INGEST] Attempt %d for batch %d", attempt + 1, i + 1)
                    vectordb.add_texts(texts=text_batch, metadatas=meta_batch)

                    # time.sleep(0.4)  # RATE LIMIT
                    START = time.time()
                    vectordb.add_texts(texts=text_batch, metadatas=meta_batch)
                    ELAPSED = time.time() - START

                    # only sleep if API was too fast
                    if ELAPSED < 0.5:
                        time.sleep(0.5 - ELAPSED)
                    batch_time = time.time() - batch_start
                    logger.info("[INGEST] Batch %d/%d done (%d chunks) in %.3fs",
                                i + 1, len(batches), len(text_batch), batch_time)
                    success = True
                    break

                except Exception as e:
                    wait = 2 ** attempt
                    logger.warning("[INGEST] Retry batch %d in %ds due to: %s", i + 1, wait, str(e))
                    logger.debug("[INGEST] Traceback for batch %d attempt %d: %s", i + 1, attempt + 1, repr(e))
                    time.sleep(wait)

            if not success:
                err_msg = f"Failed batch {i+1} for file {file_path.name}"
                logger.error("[INGEST] %s", err_msg)
                raise Exception(err_msg)

        total_time = time.time() - start_ts
        logger.info("[INGEST] Completed ingestion: %s (%d chunks) in %.3fs", file_path.name, len(chunks), total_time)

    except Exception as e:
        # log full exception with stack trace
        logger.exception("[INGEST] Failed ingestion for %s: %s", file_path.name if file_path is not None else "unknown", str(e))
        raise
